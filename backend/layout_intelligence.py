"""
layout_intelligence.py — Layout Intelligence Engine

Classifies each slide into a semantic layout type based on:
- Layout name from the template
- Number and types of placeholders
- Placeholder positions and sizes

Layout types: title, agenda, section, content, two_column,
              comparison, timeline, kpi, table, closing
"""
from pptx.util import Emu
import re

# ---------------------------------------------------------------------------
# Layout type definitions
# ---------------------------------------------------------------------------

LAYOUT_TYPES = [
    "title", "agenda", "section", "content",
    "two_column", "comparison", "timeline", "kpi", "table", "closing",
]

# Rules for content per layout type
LAYOUT_RULES = {
    "title": {
        "title_max_chars": 60,
        "body_max_chars":  120,
        "max_bullets":     1,
        "description":     "Cover slide — short title + tagline only",
    },
    "agenda": {
        "title_max_chars": 50,
        "body_max_chars":  300,
        "max_bullets":     6,
        "description":     "Agenda/table of contents — numbered list of sections",
    },
    "section": {
        "title_max_chars": 60,
        "body_max_chars":  100,
        "max_bullets":     1,
        "description":     "Section divider — section name + optional subtitle",
    },
    "content": {
        "title_max_chars": 60,
        "body_max_chars":  380,
        "max_bullets":     5,
        "description":     "Standard content slide — title + 3-5 bullet points",
    },
    "two_column": {
        "title_max_chars": 60,
        "body_max_chars":  200,   # per column
        "max_bullets":     4,     # per column
        "description":     "Two-column layout — left and right content areas",
    },
    "comparison": {
        "title_max_chars": 60,
        "body_max_chars":  200,
        "max_bullets":     4,
        "description":     "Comparison slide — two sides with labels",
    },
    "timeline": {
        "title_max_chars": 60,
        "body_max_chars":  300,
        "max_bullets":     5,
        "description":     "Timeline — sequential events or milestones",
    },
    "kpi": {
        "title_max_chars": 60,
        "body_max_chars":  250,
        "max_bullets":     4,
        "description":     "KPI/metrics slide — key numbers and data points",
    },
    "table": {
        "title_max_chars": 60,
        "body_max_chars":  400,
        "max_bullets":     8,
        "description":     "Table slide — structured rows of data",
    },
    "closing": {
        "title_max_chars": 60,
        "body_max_chars":  150,
        "max_bullets":     2,
        "description":     "Closing/thank you slide — brief closing message",
    },
}

# ---------------------------------------------------------------------------
# Classification logic
# ---------------------------------------------------------------------------

_TITLE_KEYWORDS    = ["title", "cover", "opening", "intro"]
_AGENDA_KEYWORDS   = ["agenda", "contents", "outline", "overview", "table of"]
_SECTION_KEYWORDS  = ["section", "chapter", "divider", "header", "part"]
_CLOSING_KEYWORDS  = ["thank", "closing", "end", "contact", "questions", "q&a", "goodbye"]
_KPI_KEYWORDS      = ["kpi", "metric", "data", "stat", "number", "figure", "chart", "graph"]
_TIMELINE_KEYWORDS = ["timeline", "roadmap", "schedule", "milestone", "phase", "journey"]
_TABLE_KEYWORDS    = ["table", "comparison table", "matrix", "grid"]
_COMPARISON_KEYWORDS = ["comparison", "compare", "vs", "versus", "pros", "cons"]
_TWO_COL_KEYWORDS  = ["two content", "two column", "2 column", "split", "side by side"]


def _match_keywords(text: str, keywords: list[str]) -> bool:
    t = text.lower()
    return any(k in t for k in keywords)


def classify_slide(slide_info: dict) -> str:
    """
    Classify a slide into one of the LAYOUT_TYPES.

    slide_info keys: layout (str), title_hint (str), placeholders (list), has_body (bool)
    """
    layout = slide_info.get("layout", "").lower()
    title  = slide_info.get("title_hint", "").lower()
    phs    = slide_info.get("placeholders", [])
    combined = f"{layout} {title}"

    # Count body placeholders
    body_phs = [p for p in phs if p.get("is_body")]
    n_body   = len(body_phs)

    # Closing
    if _match_keywords(combined, _CLOSING_KEYWORDS):
        return "closing"

    # Title slide (layout name is literally "title slide" or only 1 placeholder total)
    if _match_keywords(layout, _TITLE_KEYWORDS) and n_body <= 1:
        if not _match_keywords(layout, _AGENDA_KEYWORDS):
            return "title"

    # Agenda
    if _match_keywords(combined, _AGENDA_KEYWORDS):
        return "agenda"

    # Section header
    if _match_keywords(layout, _SECTION_KEYWORDS) and n_body <= 1:
        return "section"

    # Two-column
    if _match_keywords(combined, _TWO_COL_KEYWORDS) or n_body >= 2:
        # Check if body placeholders are side-by-side (similar height, different left pos)
        if n_body >= 2:
            lefts = [p.get("left", 0) for p in body_phs[:2]]
            if abs(lefts[0] - lefts[1]) > 1_000_000:  # > ~1 inch apart
                return "two_column"

    # Comparison
    if _match_keywords(combined, _COMPARISON_KEYWORDS):
        return "comparison"

    # Timeline
    if _match_keywords(combined, _TIMELINE_KEYWORDS):
        return "timeline"

    # KPI / data
    if _match_keywords(combined, _KPI_KEYWORDS):
        return "kpi"

    # Table
    if _match_keywords(combined, _TABLE_KEYWORDS):
        return "table"

    # Default: standard content slide
    return "content"


def enrich_slides(parsed_slides: list[dict]) -> list[dict]:
    """
    Add layout_type and layout_rules to each parsed slide descriptor.
    Returns enriched list (modifies in place and returns).
    """
    for slide in parsed_slides:
        lt = classify_slide(slide)
        slide["layout_type"] = lt
        slide["layout_rules"] = LAYOUT_RULES[lt]
    return parsed_slides


def slides_to_prompt_str(slides: list[dict]) -> str:
    """
    Build the LLM-readable slide structure string with layout intelligence.
    Includes layout_type, capacity, and rules so agents know exactly what to write.
    """
    lines = []
    for s in slides:
        lt    = s.get("layout_type", "content")
        rules = s.get("layout_rules", LAYOUT_RULES["content"])
        body_phs = [p for p in s.get("placeholders", []) if p.get("is_body")]
        body_cap = min((p["capacity"] for p in body_phs), default=rules["body_max_chars"])
        # Use the smaller of measured capacity and rule limit
        effective_cap = min(body_cap, rules["body_max_chars"])
        title = s.get("title_hint") or "(no title)"
        has_body = s.get("has_body", False)
        body_info = f"fillable | capacity={effective_cap} chars | max_bullets={rules['max_bullets']}" \
                    if has_body else "no-content-area"

        lines.append(
            f"Slide {s['slide_number']} | Layout: {s['layout']} | Type: {lt} | "
            f"Title: {title} | Body: {body_info}"
        )
        for p in body_phs:
            if p.get("hint"):
                lines.append(f"  Placeholder idx={p['idx']} hint: {p['hint']}")
    return "\n".join(lines)
