"""
ppt_creator.py — PPT Creator Agent

Core insight: this template uses plain text boxes (no PowerPoint placeholders).
The reliable approach is to scan the template ONCE at upload time, record the
exact shape name for each slide's title and body box, then use those names
at write time — no heuristics, no guessing.

Shape selection strategy (in priority order):
1. If the template has PowerPoint placeholders → use idx 0 (title) and idx 1 (body)
2. If no placeholders → scan all text shapes, pick:
   - Title: the shape whose existing text best matches known title patterns
             OR the topmost shape with the largest font
   - Body:  the shape with the most text content that isn't the title
             AND has area > 1/4 of the slide area (filters out small labels)
"""
from pptx import Presentation
from pptx.util import Pt
from copy import deepcopy
import os
import uuid
import re
import logging

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Markdown stripper
# ---------------------------------------------------------------------------

def _strip_markdown(text: str) -> str:
    text = re.sub(r'\*{1,3}([^*]+)\*{1,3}', r'\1', text)
    text = re.sub(r'_{1,2}([^_]+)_{1,2}', r'\1', text)
    text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'`([^`]+)`', r'\1', text)
    text = re.sub(r'^\s*[-*•]\s+', '- ', text, flags=re.MULTILINE)
    return text.strip()

# ---------------------------------------------------------------------------
# Template scanner — called once at upload time
# ---------------------------------------------------------------------------

# Known title-like template labels (lowercase)
_TITLE_LABELS = {
    "presentation title", "section title", "title", "thank you",
    "agenda", "data & insights", "data and insights",
}

# Patterns that are purely decorative (skip these)
_DECORATIVE = re.compile(
    r'^(\d{1,2}|[↑↓→←★•\|✓✗×÷]|[A-Z]{1,2}\d?)$'
)


def _get_slide_dims(prs: Presentation):
    try:
        return prs.slide_width, prs.slide_height
    except Exception:
        return 9144000, 6858000  # 10×7.5 inch default


def _max_font_pt(shape) -> float:
    """Return the largest font size (in points) found in a shape's text frame."""
    best = 0.0
    if not shape.has_text_frame:
        return best
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                sz = run.font.size
                if sz and sz > 0:
                    pt = sz / 12700.0
                    if pt > best:
                        best = pt
            except Exception:
                pass
    return best


def scan_template(template_path: str) -> list[dict]:
    """
    Scan a PPTX template and return a per-slide map of which shape names
    to use for title and body content.

    Returns list of dicts:
      { slide_number, title_shape_name, body_shape_name,
        title_hint, body_hint, has_placeholder }
    """
    prs = Presentation(template_path)
    sw, sh = _get_slide_dims(prs)
    slide_area = sw * sh
    result = []

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        entry = {
            "slide_number":      slide_num,
            "title_shape_name":  None,
            "body_shape_name":   None,
            "title_hint":        "",
            "body_hint":         "",
            "has_placeholder":   False,
        }

        # ── Pass 1: try PowerPoint placeholders ──────────────────────────
        title_ph = None
        body_ph  = None
        for shape in slide.shapes:
            if not shape.is_placeholder or not shape.has_text_frame:
                continue
            try:
                idx = shape.placeholder_format.idx
                pt  = shape.placeholder_format.type
            except Exception:
                continue
            if pt in {1, 13} or idx == 0:
                if title_ph is None:
                    title_ph = shape
            elif pt not in {3, 4, 10, 11, 12, 14} and idx >= 1:
                if body_ph is None:
                    body_ph = shape

        if title_ph or body_ph:
            entry["has_placeholder"] = True
            if title_ph:
                entry["title_shape_name"] = title_ph.name
                entry["title_hint"]       = title_ph.text_frame.text.strip()[:80]
            if body_ph:
                entry["body_shape_name"] = body_ph.name
                entry["body_hint"]       = body_ph.text_frame.text.strip()[:80]
            result.append(entry)
            continue

        # ── Pass 2: plain text boxes — score each shape ───────────────────
        # Collect all non-empty, non-decorative text shapes
        shapes_info = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if _DECORATIVE.match(text):
                continue
            try:
                area = shape.width * shape.height
                top  = shape.top
                left = shape.left
            except Exception:
                area, top, left = 0, 0, 0

            font_pt = _max_font_pt(shape)
            shapes_info.append({
                "shape":    shape,
                "text":     text,
                "area":     area,
                "top":      top,
                "left":     left,
                "font_pt":  font_pt,
                "is_title_label": text.lower() in _TITLE_LABELS,
            })

        if not shapes_info:
            result.append(entry)
            continue

        # Score each shape for "title likelihood"
        # Higher score = more likely to be the title
        def title_score(info):
            score = 0
            # Known title label → strong signal
            if info["is_title_label"]:
                score += 1000
            # Large font → likely title
            score += info["font_pt"] * 10
            # Near top of slide → likely title
            score += max(0, 500 - info["top"] / 10000)
            # Very large area alone doesn't make it a title
            return score

        shapes_info.sort(key=title_score, reverse=True)
        title_info = shapes_info[0]
        entry["title_shape_name"] = title_info["shape"].name
        entry["title_hint"]       = title_info["text"][:80]

        # Body: largest area shape that is NOT the title
        # AND has area > 5% of slide area (filters out small labels)
        min_body_area = slide_area * 0.05
        body_candidates = [
            s for s in shapes_info[1:]
            if s["area"] >= min_body_area
        ]

        if body_candidates:
            # Pick the one with the most text content (template placeholder text)
            body_info = max(body_candidates, key=lambda s: len(s["text"]))
            entry["body_shape_name"] = body_info["shape"].name
            entry["body_hint"]       = body_info["text"][:80]

        result.append(entry)

    return result

# ---------------------------------------------------------------------------
# Text writer — python-pptx native API
# ---------------------------------------------------------------------------

def _write_text(shape, lines: list[str]):
    """
    Write lines into a shape's text frame.
    Preserves font size and color from the original template.
    """
    if not lines or not shape.has_text_frame:
        return

    tf = shape.text_frame
    tf.word_wrap = True

    # Capture font properties from the first existing run
    ref_size  = None
    ref_color = None
    for para in tf.paragraphs:
        for run in para.runs:
            try:
                if run.font.size and run.font.size > 0:
                    ref_size = run.font.size
                if run.font.color and run.font.color.type:
                    ref_color = run.font.color.rgb
            except Exception:
                pass
            if ref_size:
                break
        if ref_size:
            break

    # Clear all paragraphs down to one
    while len(tf.paragraphs) > 1:
        last_p = tf.paragraphs[-1]._p
        last_p.getparent().remove(last_p)

    # Write first line into the existing first paragraph
    first_para = tf.paragraphs[0]
    for run in list(first_para.runs):
        run._r.getparent().remove(run._r)

    run0 = first_para.add_run()
    run0.text = lines[0]
    _apply_font(run0, ref_size, ref_color)

    # Add remaining lines as new paragraphs
    for line in lines[1:]:
        new_para = tf.add_paragraph()
        new_run  = new_para.add_run()
        new_run.text = line
        _apply_font(new_run, ref_size, ref_color)

    logger.info(f"    Wrote {len(lines)} line(s) to '{shape.name}': {repr(lines[0][:50])}")


def _apply_font(run, size, color):
    try:
        if size:
            run.font.size = size
        if color:
            run.font.color.rgb = color
    except Exception:
        pass

# ---------------------------------------------------------------------------
# Slide cloning
# ---------------------------------------------------------------------------

def _clone_slide(prs: Presentation, src_idx: int):
    src   = prs.slides[src_idx]
    new_s = prs.slides.add_slide(src.slide_layout)
    sp    = new_s.shapes._spTree
    # Remove all auto-generated shapes from the new slide
    for el in list(sp)[2:]:  # keep spTree header elements (nvGrpSpPr, grpSpPr)
        sp.remove(el)
    # Deep-copy every shape element from the source slide
    for el in list(src.shapes._spTree)[2:]:
        sp.append(deepcopy(el))
    return new_s


def _best_clone_source(prs: Presentation) -> int:
    if len(prs.slides) >= 3:
        return 2
    return max(0, len(prs.slides) - 2)

# ---------------------------------------------------------------------------
# Debug helper
# ---------------------------------------------------------------------------

def debug_template(template_path: str) -> list[dict]:
    prs = Presentation(template_path)
    result = []
    for i, slide in enumerate(prs.slides):
        slide_info = {"slide": i + 1, "layout": slide.slide_layout.name, "shapes": []}
        for shape in slide.shapes:
            info = {
                "name":           shape.name,
                "is_placeholder": shape.is_placeholder,
                "has_text_frame": shape.has_text_frame,
                "ph_idx":         shape.placeholder_format.idx if shape.is_placeholder else None,
                "ph_type":        shape.placeholder_format.type if shape.is_placeholder else None,
                "is_title":       False,
                "is_body":        False,
                "current_text":   shape.text_frame.text[:60] if shape.has_text_frame else "",
            }
            slide_info["shapes"].append(info)
        result.append(slide_info)
    return result

# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def create_ppt_from_template(
    template_path: str,
    slide_mappings: list[dict],
    output_dir: str = "/tmp",
    shape_map: list[dict] | None = None,
) -> str:
    """
    Fill a PPTX template with AI-mapped content.

    shape_map: pre-computed from scan_template() at upload time.
               If None, scan_template() is called here (slower but safe).
    """
    prs = Presentation(template_path)
    content_map = {int(m["slide_number"]): m for m in slide_mappings}

    # Build shape_map if not provided
    if shape_map is None:
        shape_map = scan_template(template_path)

    # Index shape_map by slide_number
    shape_index = {s["slide_number"]: s for s in shape_map}

    logger.info(f"create_ppt: {len(prs.slides)} slides, "
                f"mappings: {sorted(content_map.keys())}")

    if not content_map:
        os.makedirs(output_dir, exist_ok=True)
        out = os.path.join(output_dir, f"slideai_{uuid.uuid4().hex[:8]}.pptx")
        prs.save(out)
        return out

    # Clone extra slides if needed, and extend shape_index for cloned slides
    max_num = max(content_map.keys())
    clone_src_idx = _best_clone_source(prs)
    clone_src_num = clone_src_idx + 1  # 1-based slide number of the clone source
    while len(prs.slides) < max_num:
        new_slide_num = len(prs.slides) + 1
        _clone_slide(prs, clone_src_idx)
        # Inherit the shape map from the clone source so content gets written
        if clone_src_num in shape_index and new_slide_num not in shape_index:
            shape_index[new_slide_num] = dict(shape_index[clone_src_num])
            shape_index[new_slide_num]["slide_number"] = new_slide_num

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        mapping   = content_map.get(slide_num)
        if not mapping:
            continue

        title_text    = _strip_markdown(str(mapping.get("slide_title",      "")).strip())
        content_text  = _strip_markdown(str(mapping.get("suggested_content", "")).strip())
        content_lines = [l.strip() for l in content_text.splitlines() if l.strip()]

        logger.info(f"Slide {slide_num}: title='{title_text[:50]}', lines={len(content_lines)}")

        # Get the pre-computed shape names for this slide
        smap = shape_index.get(slide_num, {})
        title_name = smap.get("title_shape_name")
        body_name  = smap.get("body_shape_name")

        logger.info(f"  target shapes: title='{title_name}', body='{body_name}'")

        # Build a name→shape lookup for this slide
        shape_by_name = {s.name: s for s in slide.shapes if s.has_text_frame}

        title_shape = shape_by_name.get(title_name) if title_name else None
        body_shape  = shape_by_name.get(body_name)  if body_name  else None

        if title_shape and title_text:
            _write_text(title_shape, [title_text])
        elif not title_shape:
            logger.warning(f"  Slide {slide_num}: title shape '{title_name}' not found")

        if body_shape and content_lines:
            _write_text(body_shape, content_lines)
        elif not body_shape:
            logger.warning(f"  Slide {slide_num}: body shape '{body_name}' not found")

    os.makedirs(output_dir, exist_ok=True)
    out_path = os.path.join(output_dir, f"slideai_{uuid.uuid4().hex[:8]}.pptx")
    prs.save(out_path)
    logger.info(f"Saved: {out_path}")
    return out_path
