from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from dotenv import load_dotenv
import os, uuid, json, shutil, asyncio, sqlite3, time, threading
from concurrent.futures import ProcessPoolExecutor

import logging
logging.basicConfig(level=logging.INFO)

load_dotenv()

from ppt_parser import parse_ppt, slides_to_json
from layout_intelligence import enrich_slides, slides_to_prompt_str
from agent import store_template, run_agent_stream, get_template, refine_slide, session_store
from ppt_creator import create_ppt_from_template, debug_template, scan_template
from ppt_preview import pptx_to_images

app = FastAPI(title="PPT Chatbot API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

NVIDIA_API_KEY  = os.getenv("NVIDIA_API_KEY", "")
NVIDIA_BASE_URL = os.getenv("NVIDIA_BASE_URL", "https://integrate.api.nvidia.com/v1")
NVIDIA_MODEL    = os.getenv("NVIDIA_MODEL", "meta/llama-3.3-70b-instruct")

MAX_UPLOAD_BYTES = 50 * 1024 * 1024   # 50 MB

UPLOAD_DIR  = "/tmp/slideai_uploads"
PREVIEW_DIR = "/tmp/slideai_previews"
OUTPUT_DIR  = "/tmp/slideai_outputs"
DB_PATH     = "/tmp/slideai_sessions.db"

for d in [UPLOAD_DIR, PREVIEW_DIR, OUTPUT_DIR]:
    os.makedirs(d, exist_ok=True)

_preview_executor = ProcessPoolExecutor(max_workers=2)

# ---------------------------------------------------------------------------
# SQLite session persistence
# ---------------------------------------------------------------------------

def _init_db():
    conn = sqlite3.connect(DB_PATH)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS sessions (
            session_id   TEXT PRIMARY KEY,
            slide_struct TEXT,
            template_path TEXT,
            char_limits  TEXT,
            shape_map    TEXT,
            created_at   REAL
        )
    """)
    conn.commit()
    conn.close()


def _save_session(session_id: str, slide_struct: str, template_path: str,
                  char_limits: dict, shape_map: list):
    conn = sqlite3.connect(DB_PATH)
    conn.execute("""
        INSERT OR REPLACE INTO sessions VALUES (?,?,?,?,?,?)
    """, (session_id, slide_struct, template_path,
          json.dumps(char_limits), json.dumps(shape_map), time.time()))
    conn.commit()
    conn.close()


def _load_session(session_id: str) -> dict | None:
    conn = sqlite3.connect(DB_PATH)
    row  = conn.execute(
        "SELECT slide_struct, template_path, char_limits, shape_map FROM sessions WHERE session_id=?",
        (session_id,)
    ).fetchone()
    conn.close()
    if not row:
        return None
    return {
        "slide_struct":   row[0],
        "template_path":  row[1],
        "char_limits":    json.loads(row[2]),
        "shape_map":      json.loads(row[3]),
    }

# ---------------------------------------------------------------------------
# Temp file cleanup (files older than 24 h)
# ---------------------------------------------------------------------------

def _cleanup_old_files(directory: str, max_age_hours: int = 24):
    cutoff = time.time() - (max_age_hours * 3600)
    for root, dirs, files in os.walk(directory, topdown=False):
        for f in files:
            path = os.path.join(root, f)
            try:
                if os.path.getmtime(path) < cutoff:
                    os.remove(path)
            except OSError:
                pass
        for d in dirs:
            dpath = os.path.join(root, d)
            try:
                os.rmdir(dpath)
            except OSError:
                pass


def _run_cleanup():
    for d in [UPLOAD_DIR, PREVIEW_DIR, OUTPUT_DIR]:
        _cleanup_old_files(d, max_age_hours=24)
    logging.info("[cleanup] Temp file cleanup complete")

# ---------------------------------------------------------------------------
# In-memory stores (populated from DB on demand)
# ---------------------------------------------------------------------------

template_file_store: dict[str, str]  = {}
shape_map_store:     dict[str, list] = {}
char_limit_store:    dict[str, dict] = {}


def _restore_session(session_id: str) -> bool:
    """Restore session from SQLite if not in memory. Returns True if found."""
    if session_id in template_file_store:
        return True
    data = _load_session(session_id)
    if not data:
        return False
    template_path = data["template_path"]
    if not os.path.exists(template_path):
        return False
    store_template(session_id, data["slide_struct"])
    template_file_store[session_id] = template_path
    shape_map_store[session_id]     = data["shape_map"]
    char_limit_store[session_id]    = {int(k): v for k, v in data["char_limits"].items()}
    logging.info(f"[restore] Session {session_id} restored from DB")
    return True

# ---------------------------------------------------------------------------
# Startup
# ---------------------------------------------------------------------------

@app.on_event("startup")
async def startup():
    _init_db()
    # Run cleanup in background thread so startup is non-blocking
    threading.Thread(target=_run_cleanup, daemon=True).start()
    logging.info(f"[startup] model={NVIDIA_MODEL} base_url={NVIDIA_BASE_URL}")

app.mount("/previews", StaticFiles(directory=PREVIEW_DIR), name="previews")

# ---------------------------------------------------------------------------
# Health check
# ---------------------------------------------------------------------------

@app.get("/health")
def health():
    return {
        "status":               "ok",
        "sessions_in_memory":   len(session_store),
        "nvidia_key_set":       bool(NVIDIA_API_KEY),
        "libreoffice_available": shutil.which("libreoffice") is not None,
        "db_path":              DB_PATH,
    }

# ---------------------------------------------------------------------------
# Session
# ---------------------------------------------------------------------------

@app.get("/new-session")
def new_session():
    return {"session_id": str(uuid.uuid4())}

# ---------------------------------------------------------------------------
# Upload PPT Template
# ---------------------------------------------------------------------------

@app.post("/upload-template")
async def upload_template(file: UploadFile = File(...), session_id: str = Form(...)):
    if not file.filename.endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Only .pptx files are supported")

    # Read and validate file size before writing to disk
    content = await file.read()
    if len(content) > MAX_UPLOAD_BYTES:
        raise HTTPException(status_code=400, detail="File too large. Maximum size is 50 MB.")

    save_path = os.path.join(UPLOAD_DIR, f"{session_id}.pptx")
    with open(save_path, "wb") as f:
        f.write(content)

    slides          = parse_ppt(save_path)
    enriched_slides = enrich_slides(slides)
    slide_struct    = slides_to_prompt_str(enriched_slides)

    store_template(session_id, slide_struct)
    template_file_store[session_id] = save_path

    shape_map = scan_template(save_path)
    shape_map_store[session_id] = shape_map

    char_limits = {
        s["slide_number"]: s.get("layout_rules", {}).get("body_max_chars", 400)
        for s in enriched_slides
    }
    char_limit_store[session_id] = char_limits

    # Persist to SQLite so session survives server restarts
    _save_session(session_id, slide_struct, save_path, char_limits, shape_map)

    logging.info(f"[upload] session={session_id} slides={len(slides)} "
                 f"layout_types={[s.get('layout_type') for s in enriched_slides]}")

    return {
        "message":         "Template uploaded and parsed successfully",
        "session_id":      session_id,
        "slides_detected": len(slides),
        "slide_structure": slide_struct,
        "layout_summary":  [
            {"slide": s["slide_number"], "layout_type": s.get("layout_type", "content")}
            for s in enriched_slides
        ],
    }

# ---------------------------------------------------------------------------
# Process Raw Data (streaming SSE)
# ---------------------------------------------------------------------------

class ProcessRequest(BaseModel):
    session_id: str
    raw_data:   str
    tone:       str = "General"


@app.post("/process-data")
async def process_data(request: ProcessRequest):
    if not NVIDIA_API_KEY:
        raise HTTPException(status_code=500, detail="NVIDIA_API_KEY not set on server")

    if not _restore_session(request.session_id):
        raise HTTPException(status_code=400, detail="No PPT template found. Please upload your PPT first.")

    async def event_stream():
        async for chunk in run_agent_stream(
            session_id=request.session_id,
            raw_data=request.raw_data,
            openai_api_key=NVIDIA_API_KEY,
            base_url=NVIDIA_BASE_URL,
            model=NVIDIA_MODEL,
            tone=request.tone,
        ):
            yield f"data: {json.dumps(chunk)}\n\n"

    return StreamingResponse(event_stream(), media_type="text/event-stream")

# ---------------------------------------------------------------------------
# Create PPT
# ---------------------------------------------------------------------------

class CreatePPTRequest(BaseModel):
    session_id: str
    slides:     list[dict]


@app.post("/create-ppt")
async def create_ppt(request: CreatePPTRequest):
    _restore_session(request.session_id)
    template_path = template_file_store.get(request.session_id)
    if not template_path or not os.path.exists(template_path):
        raise HTTPException(status_code=400, detail="Original template not found. Please re-upload your PPT.")

    try:
        output_path = create_ppt_from_template(
            template_path=template_path,
            slide_mappings=request.slides,
            output_dir=OUTPUT_DIR,
            shape_map=shape_map_store.get(request.session_id),
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PPT creation failed: {str(e)}")

    return FileResponse(
        path=output_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="SlideAI_Output.pptx",
        background=None,
    )

# ---------------------------------------------------------------------------
# Preview PPT — non-blocking LibreOffice via executor
# ---------------------------------------------------------------------------

class PreviewRequest(BaseModel):
    session_id: str
    slides:     list[dict]


def _run_preview_sync(template_path: str, slides: list, shape_map: list, preview_dir: str):
    filled_path = create_ppt_from_template(
        template_path=template_path,
        slide_mappings=slides,
        output_dir=OUTPUT_DIR,
        shape_map=shape_map,
    )
    return pptx_to_images(filled_path, output_dir=preview_dir)


@app.post("/preview-ppt")
async def preview_ppt(request: PreviewRequest):
    _restore_session(request.session_id)
    template_path = template_file_store.get(request.session_id)
    if not template_path or not os.path.exists(template_path):
        raise HTTPException(status_code=400, detail="Template not found. Please re-upload your PPT.")

    try:
        loop        = asyncio.get_event_loop()
        image_paths = await loop.run_in_executor(
            _preview_executor,
            _run_preview_sync,
            template_path,
            request.slides,
            shape_map_store.get(request.session_id),
            PREVIEW_DIR,
        )
        image_urls = [f"/previews/{os.path.relpath(p, PREVIEW_DIR)}" for p in image_paths]
        return {"slides": image_urls, "total": len(image_urls)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Preview generation failed: {str(e)}")

# ---------------------------------------------------------------------------
# Refine Slide
# ---------------------------------------------------------------------------

class RefineSlideRequest(BaseModel):
    session_id:      str
    slide_number:    int
    slide_title:     str
    current_content: str
    instruction:     str
    chat_history:    list = []


@app.post("/refine-slide")
async def refine_slide_endpoint(request: RefineSlideRequest):
    if not _restore_session(request.session_id):
        raise HTTPException(
            status_code=400,
            detail=f"Session '{request.session_id}' not found. Please upload a PPT template first.",
        )

    char_limit = char_limit_store.get(request.session_id, {}).get(request.slide_number, 400)

    try:
        refined_content = await refine_slide(
            session_id=request.session_id,
            slide_number=request.slide_number,
            slide_title=request.slide_title,
            current_content=request.current_content,
            instruction=request.instruction,
            chat_history=request.chat_history,
            openai_api_key=NVIDIA_API_KEY,
            base_url=NVIDIA_BASE_URL,
            model=NVIDIA_MODEL,
            char_limit=char_limit,
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"LLM refinement failed: {str(e)}")

    return {"refined_content": refined_content, "char_limit": char_limit}

# ---------------------------------------------------------------------------
# Debug endpoints
# ---------------------------------------------------------------------------

@app.get("/debug-template/{session_id}")
def debug_template_endpoint(session_id: str):
    _restore_session(session_id)
    template_path = template_file_store.get(session_id)
    if not template_path or not os.path.exists(template_path):
        raise HTTPException(status_code=404, detail="Template not found")
    return {"slides": debug_template(template_path)}


@app.post("/debug-write/{session_id}")
def debug_write_endpoint(session_id: str):
    _restore_session(session_id)
    template_path = template_file_store.get(session_id)
    if not template_path or not os.path.exists(template_path):
        raise HTTPException(status_code=404, detail="Template not found")

    from pptx import Presentation
    prs = Presentation(template_path)
    test_mappings = [
        {
            "slide_number":      i + 1,
            "slide_title":       f"TEST TITLE SLIDE {i+1}",
            "suggested_content": "- Test bullet point one\n- Test bullet point two\n- Test bullet point three",
            "reason":            "debug test",
        }
        for i, _ in enumerate(prs.slides)
    ]
    try:
        out = create_ppt_from_template(
            template_path=template_path,
            slide_mappings=test_mappings,
            output_dir=OUTPUT_DIR,
            shape_map=shape_map_store.get(session_id),
        )
        return {"output_path": out, "mappings_applied": len(test_mappings)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
