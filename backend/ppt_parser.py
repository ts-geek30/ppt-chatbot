"""
ppt_parser.py — Parse a PPTX template and extract a rich, structured
description of every slide's placeholders so the LLM and creator can
work with precise placeholder-level information.
"""
from pptx import Presentation
from pptx.util import Emu
import json
import re

# ---------------------------------------------------------------------------
# Placeholder classification (shared constants)
# ---------------------------------------------------------------------------
TITLE_TYPES  = {1, 13}          # TITLE, CENTER_TITLE
BODY_TYPES   = {2, 7, 15}       # BODY, OBJECT, VERTICAL_BODY
SKIP_TYPES   = {3, 4, 10, 11, 12, 14}  # DATE, FOOTER, SLIDE_NUMBER, etc.


def ph_type(shape):
    return shape.placeholder_format.type if shape.is_placeholder else None


def ph_idx(shape):
    return shape.placeholder_format.idx if shape.is_placeholder else None


def is_title(shape) -> bool:
    pt = ph_type(shape)
    if pt in TITLE_TYPES:
        return True
    if shape.is_placeholder and ph_idx(shape) == 0:
        return True
    return False


def is_fillable_body(shape) -> bool:
    """True if this shape is a writable content placeholder."""
    if not shape.has_text_frame or not shape.is_placeholder:
        return False
    if is_title(shape):
        return False
    pt  = ph_type(shape)
    idx = ph_idx(shape)
    if pt in SKIP_TYPES:
        return False
    if idx is not None and idx >= 1:
        return True
    if pt in BODY_TYPES:
        return True
    return False


def _read_font_pt(shape, fallback: float = 16.0) -> float:
    """
    Read the actual font size from a shape's text frame.
    Falls back to `fallback` if no explicit size is set.
    This fixes the previous hardcoded 18pt default which underestimated capacity by ~20%.
    """
    if not shape.has_text_frame:
        return fallback
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                sz = run.font.size
                if sz and sz > 0:
                    pt = sz / 12700.0
                    # Sanity check: ignore unrealistic sizes
                    if 6.0 <= pt <= 72.0:
                        return pt
            except Exception:
                pass
    return fallback


def _emu_to_chars(width_emu: int, height_emu: int, font_pt: float) -> dict:
    """
    Estimate how many characters fit in a text box given its actual font size.
    Returns capacity and spatial hints.
    """
    width_pt  = width_emu  / 12700.0
    height_pt = height_emu / 12700.0

    # 0.55 pts per char width (slightly tighter than old 0.6 — more accurate for common fonts)
    chars_per_line = max(15, int(width_pt / (font_pt * 0.55)))
    # 1.4x font size for line height (accounts for typical line spacing)
    lines = max(1, int(height_pt / (font_pt * 1.4)))

    capacity = chars_per_line * lines

    aspect = width_pt / height_pt if height_pt > 0 else 1
    if aspect > 2:
        shape_type = "wide/short"
    elif aspect < 0.5:
        shape_type = "narrow/tall"
    else:
        shape_type = "balanced"

    return {"capacity": capacity, "shape_type": shape_type, "font_pt": round(font_pt, 1)}


def parse_ppt(file_path: str) -> list[dict]:
    """
    Parse a PPTX and return slide descriptors.
    Works for both placeholder-based and plain text box templates.
    """
    prs = Presentation(file_path)
    slides = []

    _title_hints_lower = {
        "presentation title", "section title", "title", "thank you",
        "agenda", "data & insights",
    }
    _skip_re = re.compile(r'^(\d{1,2}|[↑↓→←★•\|])')

    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
        placeholders = []
        title_hint   = ""

        # Collect all text shapes sorted by vertical position
        text_shapes = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if _skip_re.match(text):
                continue
            try:
                top = shape.top
            except Exception:
                top = 0
            text_shapes.append((top, shape, text))

        text_shapes.sort(key=lambda x: x[0])

        for top, shape, text in text_shapes:
            hint = text[:100]
            try:
                # Read actual font size from the shape — fixes the 18pt hardcode bug
                font_pt = _read_font_pt(shape, fallback=16.0)
                stats   = _emu_to_chars(shape.width, shape.height, font_pt)
                cap     = stats["capacity"]
                stype   = stats["shape_type"]
            except Exception:
                cap   = 300
                stype = "balanced"

            is_t = text.lower() in _title_hints_lower
            is_b = not is_t

            if is_t and not title_hint:
                title_hint = hint

            placeholders.append({
                "idx":        None,
                "type_code":  None,
                "hint":       hint,
                "capacity":   cap,
                "shape_type": stype,
                "is_title":   is_t,
                "is_body":    is_b,
            })

        if not title_hint and text_shapes:
            title_hint = text_shapes[0][2][:100]

        has_body = len(text_shapes) > 1

        slides.append({
            "slide_number": i + 1,
            "layout":       layout_name,
            "title_hint":   title_hint,
            "has_body":     has_body,
            "placeholders": placeholders,
        })

    return slides


def slides_to_prompt_str(slides: list[dict]) -> str:
    """
    Fallback prompt string builder (used when layout_intelligence is not available).
    Prefer layout_intelligence.slides_to_prompt_str() for enriched output.
    """
    lines = []
    for s in slides:
        body_shapes = [p for p in s["placeholders"] if p["hint"]]
        if body_shapes:
            max_cap    = max(p["capacity"] for p in body_shapes)
            best_shape = next(p for p in body_shapes if p["capacity"] == max_cap)
            stype      = best_shape.get("shape_type", "balanced")
            body_info  = f"fillable ({stype}, max ~{max_cap} chars)"
        else:
            body_info = "no-content-area"
        title = s["title_hint"] or "(no title)"
        lines.append(
            f"Slide {s['slide_number']} | Layout: {s['layout']} | Title: {title} | Body: {body_info}"
        )
        for p in body_shapes[:3]:
            lines.append(f"  Text hint: {p['hint']} (capacity ~{p['capacity']} chars, {p.get('shape_type')})")
    return "\n".join(lines)


def slides_to_json(slides: list[dict]) -> str:
    """Return the full parsed slide data as JSON."""
    return json.dumps(slides)
